"""
create_presentation.py
Generates a professional PowerPoint presentation for the VaR project.
Run: python SCRIPT/create_presentation.py
Output: VaR_Presentation.pptx (in project root)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Paths ──────────────────────────────────────────────────────────────────────
BASE_DIR  = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DATA_DIR  = os.path.join(BASE_DIR, "DATA")
OUT_FILE  = os.path.join(BASE_DIR, "VaR_Presentation.pptx")

CHART = {
    1: os.path.join(DATA_DIR, "chart_01_prices.png"),
    2: os.path.join(DATA_DIR, "chart_02_portfolio_returns.png"),
    3: os.path.join(DATA_DIR, "chart_03_histogram_var.png"),
    4: os.path.join(DATA_DIR, "chart_04_var_comparison.png"),
}

# ── Colour palette ─────────────────────────────────────────────────────────────
DARK_BLUE  = RGBColor(0x1F, 0x3B, 0x6E)   # titles / accents
MID_BLUE   = RGBColor(0x2E, 0x75, 0xB6)   # sub-headings
LIGHT_GREY = RGBColor(0xF2, 0xF2, 0xF2)   # slide backgrounds
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x1A, 0x1A, 0x1A)

# ── Helpers ────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)   # widescreen 16:9
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank layout


def add_slide():
    return prs.slides.add_slide(BLANK)


def fill_background(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color: RGBColor):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def text_box(slide, text, left, top, width, height,
             font_size=18, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def bullet_box(slide, items, left, top, width, height,
               font_size=17, color=BLACK, indent_pts=14):
    """Add a text box with bullet points (one per item)."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(6)
        run = p.runs[0]
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        # manual bullet via paragraph indent + leader character
        from pptx.util import Pt as pPt
        pf = p._pPf if hasattr(p, '_pPf') else None
        # simple prefix bullet
        p.text = "▸  " + item
        run = p.runs[0]
        run.font.size = Pt(font_size)
        run.font.color.rgb = color

    return txBox


def add_notes(slide, notes_text: str):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text


def add_image(slide, img_path, left, top, width, height=None):
    if not os.path.exists(img_path):
        print(f"  [WARNING] Image not found: {img_path}")
        return
    if height:
        slide.shapes.add_picture(img_path,
                                 Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    else:
        slide.shapes.add_picture(img_path,
                                 Inches(left), Inches(top),
                                 Inches(width))


def slide_header(slide, title, subtitle=None):
    """Dark blue header bar at the top of a content slide."""
    add_rect(slide, 0, 0, 13.33, 1.2, DARK_BLUE)
    text_box(slide, title,
             left=0.4, top=0.18, width=12.5, height=0.8,
             font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        text_box(slide, subtitle,
                 left=0.4, top=0.85, width=12.5, height=0.4,
                 font_size=14, bold=False, color=RGBColor(0xBD, 0xD7, 0xEE),
                 align=PP_ALIGN.LEFT)


def footer(slide, text="Master 1 Finance | IAE Metz | 2026"):
    add_rect(slide, 0, 7.1, 13.33, 0.4, DARK_BLUE)
    text_box(slide, text,
             left=0.3, top=7.12, width=12.7, height=0.3,
             font_size=10, bold=False, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s1 = add_slide()
fill_background(s1, DARK_BLUE)

# decorative accent bar
add_rect(s1, 0, 3.0, 13.33, 0.08, MID_BLUE)

text_box(s1,
         "Value at Risk (VaR) Analysis",
         left=1.0, top=1.5, width=11.0, height=1.2,
         font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

text_box(s1,
         "A Comparative Study: Historical, Parametric, and Monte Carlo Methods",
         left=1.0, top=2.7, width=11.0, height=0.7,
         font_size=20, bold=False, color=RGBColor(0xBD, 0xD7, 0xEE),
         align=PP_ALIGN.CENTER)

text_box(s1,
         "US Equity Portfolio  •  AAPL · JPM · JNJ · XOM · WMT",
         left=1.0, top=3.45, width=11.0, height=0.5,
         font_size=16, bold=False, color=RGBColor(0xBD, 0xD7, 0xEE),
         align=PP_ALIGN.CENTER)

text_box(s1,
         "Master 1 Finance  |  IAE Metz  |  April 10, 2026",
         left=1.0, top=5.8, width=11.0, height=0.5,
         font_size=14, bold=False, color=RGBColor(0x9D, 0xB2, 0xD8),
         align=PP_ALIGN.CENTER)

add_notes(s1,
    "Good morning / afternoon everyone. My name is [your name] and today I will present "
    "my Python project on Value at Risk analysis. The goal was to measure the daily market "
    "risk of a US equity portfolio using three different statistical methods. Let me walk "
    "you through each step."
)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Objective
# ══════════════════════════════════════════════════════════════════════════════
s2 = add_slide()
fill_background(s2, LIGHT_GREY)
slide_header(s2, "Project Objective", "Why measure risk?")
footer(s2)

bullets_s2 = [
    "Estimate the maximum expected loss of a portfolio over one trading day",
    "Confidence level: 95%  →  we accept a 5% probability of exceeding the VaR",
    "Compare three standard VaR methods: Historical, Parametric, Monte Carlo",
    "Implement the entire workflow in Python (data download → calculation → charts)",
    "Provide actionable risk insight for a real equity portfolio",
]
bullet_box(s2, bullets_s2,
           left=0.5, top=1.4, width=12.3, height=5.3,
           font_size=18)

add_notes(s2,
    "The main objective of this project is to measure risk quantitatively. "
    "Value at Risk answers a simple question: 'What is the worst loss I can expect on a typical bad day?' "
    "We chose a 95% confidence level, which is the industry standard. "
    "This means that on 95 out of 100 days, the actual loss should be smaller than our VaR estimate. "
    "The key contribution of this project is comparing three different methods to see if they agree."
)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Data & Portfolio
# ══════════════════════════════════════════════════════════════════════════════
s3 = add_slide()
fill_background(s3, LIGHT_GREY)
slide_header(s3, "Data & Portfolio Construction")
footer(s3)

# Left column — text
bullets_s3 = [
    "Data source: Yahoo Finance via the yfinance Python library",
    "Period: January 2019 – December 2024  (≈ 1,500 daily observations)",
    "Prices: Adjusted daily closing prices",
    "5 stocks — equally weighted (20% each):",
    "     AAPL  — Technology (Apple)",
    "     JPM   — Financials (JPMorgan Chase)",
    "     JNJ   — Healthcare (Johnson & Johnson)",
    "     XOM   — Energy (ExxonMobil)",
    "     WMT   — Consumer staples (Walmart)",
    "Equal weighting ensures sector diversification",
]
bullet_box(s3, bullets_s3,
           left=0.5, top=1.35, width=6.8, height=5.5,
           font_size=15)

# Right column — chart
add_image(s3, CHART[1], left=7.5, top=1.3, width=5.5, height=4.0)
text_box(s3, "Fig. 1 — Normalised stock prices (2019–2024)",
         left=7.5, top=5.35, width=5.5, height=0.4,
         font_size=10, color=RGBColor(0x70, 0x70, 0x70), align=PP_ALIGN.CENTER)

add_notes(s3,
    "We downloaded five years of daily price data from Yahoo Finance using the yfinance library in Python. "
    "The portfolio is equally weighted across five stocks from different sectors: technology, financials, "
    "healthcare, energy, and consumer staples. This diversification is intentional — it reflects a realistic "
    "retail investor portfolio. The chart on the right shows how each stock performed over the period, "
    "normalised to 100 at the start."
)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Methodology
# ══════════════════════════════════════════════════════════════════════════════
s4 = add_slide()
fill_background(s4, LIGHT_GREY)
slide_header(s4, "Methodology — Three Ways to Compute VaR")
footer(s4)

col_w = 3.9
col_top = 1.4
col_h = 5.5
gap = 0.2

for i, (col_left, title, items, note_color) in enumerate([
    (0.3, "Historical VaR", [
        "Uses actual past returns",
        "No assumptions on distribution",
        "Take the 5th percentile of observed returns",
        "Simple but backward-looking",
    ], MID_BLUE),
    (4.4, "Parametric VaR", [
        "Assumes returns follow a Normal distribution",
        "Uses mean and standard deviation",
        "Applies the z-score: 1.645 at 95%",
        "Fast — but sensitive to the normality assumption",
    ], RGBColor(0x70, 0x55, 0xA6)),
    (8.5, "Monte Carlo VaR", [
        "Simulates 10,000 random return scenarios",
        "Based on historical mean & volatility",
        "Takes the 5th percentile of simulated returns",
        "Most flexible — computationally intensive",
    ], RGBColor(0x2E, 0x86, 0x48)),
]):
    # coloured header card
    add_rect(s4, col_left, col_top, col_w, 0.7, note_color)
    text_box(s4, title,
             left=col_left + 0.1, top=col_top + 0.1,
             width=col_w - 0.2, height=0.55,
             font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    bullet_box(s4, items,
               left=col_left + 0.1, top=col_top + 0.8,
               width=col_w - 0.2, height=4.3,
               font_size=14)

add_notes(s4,
    "We computed VaR using three methods. "
    "The Historical method simply looks at what happened in the past — if we sort all daily returns and find "
    "the 5th percentile, that is our VaR. No assumptions needed. "
    "The Parametric method assumes returns are normally distributed and uses the well-known formula with "
    "mean, standard deviation, and the z-score of 1.645 for 95% confidence. "
    "Finally, Monte Carlo simulates thousands of possible futures based on the historical statistics and "
    "picks the worst 5%. It is the most powerful method but requires more computation."
)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Portfolio Returns
# ══════════════════════════════════════════════════════════════════════════════
s5 = add_slide()
fill_background(s5, LIGHT_GREY)
slide_header(s5, "Portfolio Returns Over Time")
footer(s5)

add_image(s5, CHART[2], left=0.4, top=1.3, width=8.0, height=4.5)

bullets_s5 = [
    "Daily returns range roughly from −5% to +5%",
    "Returns cluster around zero — typical for equities",
    "Volatility spikes visible in 2020 (COVID-19 crash)",
    "Portfolio appears relatively stable post-2021",
]
bullet_box(s5, bullets_s5,
           left=8.7, top=1.6, width=4.3, height=4.5,
           font_size=15)

text_box(s5, "Fig. 2 — Daily portfolio returns (2019–2024)",
         left=0.4, top=5.85, width=8.0, height=0.35,
         font_size=10, color=RGBColor(0x70, 0x70, 0x70), align=PP_ALIGN.CENTER)

add_notes(s5,
    "This chart shows the daily return of our equally weighted portfolio over the five-year period. "
    "You can see that most daily moves are small — between minus two and plus two percent. "
    "However, in early 2020, when COVID-19 hit the markets, we see very large swings. "
    "This kind of volatility clustering — calm periods followed by turbulent ones — is a well-known "
    "feature of financial data. It is one reason why measuring risk is important and non-trivial."
)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Distribution & Risk
# ══════════════════════════════════════════════════════════════════════════════
s6 = add_slide()
fill_background(s6, LIGHT_GREY)
slide_header(s6, "Return Distribution & VaR Threshold")
footer(s6)

add_image(s6, CHART[3], left=0.4, top=1.3, width=8.0, height=4.5)

bullets_s6 = [
    "Histogram of all daily portfolio returns",
    "Red dashed line = Historical VaR at −1.65%",
    "5% of returns fall to the left of this line",
    "The distribution has a slight left tail (negative skew)",
    "Not perfectly normal — justifies using multiple methods",
]
bullet_box(s6, bullets_s6,
           left=8.7, top=1.6, width=4.3, height=4.5,
           font_size=15)

text_box(s6, "Fig. 3 — Return histogram with Historical VaR threshold",
         left=0.4, top=5.85, width=8.0, height=0.35,
         font_size=10, color=RGBColor(0x70, 0x70, 0x70), align=PP_ALIGN.CENTER)

add_notes(s6,
    "This histogram shows all the daily returns of the portfolio. "
    "The red dashed vertical line represents our Historical VaR: approximately minus 1.65%. "
    "Everything to the left of that line is in the worst 5% of days. "
    "Notice that the distribution has a slightly heavier left tail than a perfect normal distribution would. "
    "This is why the Parametric VaR, which assumes normality, gives a slightly different result than the "
    "historical approach."
)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — VaR Results Comparison
# ══════════════════════════════════════════════════════════════════════════════
s7 = add_slide()
fill_background(s7, LIGHT_GREY)
slide_header(s7, "VaR Results — Comparison of Three Methods")
footer(s7)

# summary table (manual)
table_left, table_top = 0.5, 1.4
col_widths = [3.5, 3.0, 3.0]
row_height = 0.6
headers = ["Method", "VaR (1-day, 95%)", "Interpretation"]
rows = [
    ["Historical",   "−1.65%", "Based on actual past returns"],
    ["Parametric",   "−1.87%", "Assumes Normal distribution"],
    ["Monte Carlo",  "−1.86%", "10,000 simulated scenarios"],
]

# header row
x = table_left
for j, (h, w) in enumerate(zip(headers, col_widths)):
    add_rect(s7, x, table_top, w - 0.05, row_height, DARK_BLUE)
    text_box(s7, h,
             left=x + 0.1, top=table_top + 0.1,
             width=w - 0.25, height=row_height - 0.1,
             font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    x += w

row_colors = [RGBColor(0xDE, 0xEB, 0xF7), WHITE, RGBColor(0xDE, 0xEB, 0xF7)]
for i, (row, rc) in enumerate(zip(rows, row_colors)):
    y = table_top + row_height * (i + 1)
    x = table_left
    for j, (cell, w) in enumerate(zip(row, col_widths)):
        add_rect(s7, x, y, w - 0.05, row_height, rc)
        text_box(s7, cell,
                 left=x + 0.1, top=y + 0.1,
                 width=w - 0.25, height=row_height - 0.1,
                 font_size=14, bold=(i == -1), color=BLACK,
                 align=PP_ALIGN.CENTER)
        x += w

# key takeaway box
add_rect(s7, 0.5, table_top + row_height * 4 + 0.15,
         9.5 - 0.1, 0.7, RGBColor(0xFF, 0xF2, 0xCC))
text_box(s7,
         "Key insight:  All three methods converge around −1.65% to −1.87%, "
         "confirming the robustness of the risk estimate.",
         left=0.6, top=table_top + row_height * 4 + 0.2,
         width=9.2, height=0.65,
         font_size=14, bold=False, color=RGBColor(0x7F, 0x60, 0x00),
         align=PP_ALIGN.LEFT)

# chart on the right
add_image(s7, CHART[4], left=9.7, top=1.3, width=3.4, height=4.5)
text_box(s7, "Fig. 4 — VaR comparison chart",
         left=9.7, top=5.85, width=3.4, height=0.35,
         font_size=10, color=RGBColor(0x70, 0x70, 0x70), align=PP_ALIGN.CENTER)

add_notes(s7,
    "Here are the final results. The Historical method gives a VaR of minus 1.65%, "
    "the Parametric gives minus 1.87%, and the Monte Carlo gives minus 1.86%. "
    "In practical terms, this means that on any given day, we can say with 95% confidence "
    "that the portfolio will not lose more than about 1.7 to 1.9 percent of its value. "
    "The fact that all three methods give similar results is reassuring — it means our estimate "
    "is stable and not highly sensitive to the method chosen. "
    "The small difference between Historical and the other two comes from the slight non-normality "
    "of the actual return distribution."
)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Conclusion & Limitations
# ══════════════════════════════════════════════════════════════════════════════
s8 = add_slide()
fill_background(s8, LIGHT_GREY)
slide_header(s8, "Conclusion & Limitations")
footer(s8)

# two columns
text_box(s8, "Conclusions",
         left=0.4, top=1.35, width=5.8, height=0.45,
         font_size=18, bold=True, color=MID_BLUE)
bullets_conc = [
    "Successfully computed 1-day VaR at 95% using Python",
    "All three methods produce consistent estimates (~1.65–1.87%)",
    "The portfolio shows moderate daily risk — typical for a diversified equity basket",
    "Monte Carlo and Parametric agree closely, confirming the near-normality assumption",
    "Python proved effective for end-to-end financial risk analysis",
]
bullet_box(s8, bullets_conc,
           left=0.4, top=1.85, width=5.9, height=4.5,
           font_size=15)

# separator
add_rect(s8, 6.5, 1.3, 0.05, 5.5, RGBColor(0xCC, 0xCC, 0xCC))

text_box(s8, "Limitations & Extensions",
         left=6.7, top=1.35, width=6.3, height=0.45,
         font_size=18, bold=True, color=RGBColor(0xC0, 0x50, 0x28))
bullets_lim = [
    "VaR does not tell us how large losses can be beyond the threshold (tail risk)",
    "Parametric VaR assumes normality — real returns have fat tails",
    "Historical method is backward-looking and misses future regime changes",
    "Extension: Expected Shortfall (CVaR) — average of the worst 5% of losses",
    "Extension: Longer horizons (10-day VaR, e.g. for Basel III compliance)",
    "Extension: Dynamic volatility models (GARCH) for time-varying risk",
]
bullet_box(s8, bullets_lim,
           left=6.7, top=1.85, width=6.3, height=4.5,
           font_size=15)

add_notes(s8,
    "To conclude, this project achieved its goal: we successfully measured and compared the daily VaR "
    "of a diversified US equity portfolio using three methods in Python. "
    "The results are consistent across methods, which gives us confidence in the estimate. "
    "However, VaR has well-known limitations. Most importantly, it does not tell us what the loss "
    "will be on the 5% of bad days — it only tells us the threshold. "
    "A natural extension would be Expected Shortfall, also called CVaR, which measures the average "
    "loss in the worst cases. I am happy to take any questions. Thank you."
)

# ── Save ───────────────────────────────────────────────────────────────────────
prs.save(OUT_FILE)
print(f"\nDone! Presentation saved: {OUT_FILE}")
print(f"   Slides: {len(prs.slides)}")
